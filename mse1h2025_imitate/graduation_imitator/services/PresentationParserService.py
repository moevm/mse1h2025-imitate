from pptx import Presentation
from pptx.slide import Slides, Slide

from ..dto.PresentationData import PresentationData
from ..configs import presentation_config
from ..utils import presentation_utils


class PresentationParserService:
    '''Class for parsing PPTX presentations'''

    @staticmethod
    def parsePPTX(pathToFile: str) -> PresentationData:
        '''Method to parse all needed info from pptx presentation
        Args:
            pathToFile: str - path to pptx presentation file
        Returns:
            PresentationData object
        Raises:
            pptx.exc.PackageNotFoundError - if file does not exists or file with invalid format
        ''' 
        presentation = Presentation(pathToFile)
        topic = PresentationParserService.__getTopic(presentation.slides)
        goalAndTasks = PresentationParserService.__getGoalAndTasks(presentation.slides)
        author = PresentationParserService.__getAuthor(presentation.slides)
        slidesTitles = [PresentationParserService.__getSlideTitle(slide) for slide in presentation.slides]
        return PresentationData(topic, goalAndTasks, author, slidesTitles)
    
    @staticmethod
    def __getTopic(slides: Slides) -> str:
        '''Method to find topic of presentation. Just trying to get title at first slide.
        Args:
            slides: Slides - presentation slides sequence
        Returns:
            str - title or `Not found`
        '''
        if not slides:
            return 'Not found'
        return PresentationParserService.__getSlideTitle(slides[0])
    
    @staticmethod
    def __getAuthor(slides: Slides) -> str:
        '''Method to find author of presentation. 
        Gets str between two marks in text at first slide (config.LEFT_AUTHOR_MARK, config.RIGHT_AUTHOR_MARK).
        Args:
            slides: Slides - presentation slides sequence
        Returns:
            str - author or `Not found`
        '''
        if not slides:
            return 'Not found'
        frontPageText = PresentationParserService.__getSlideText(slides[0]).lower()
        leftPointer = frontPageText.find(presentation_config.LEFT_AUTHOR_MARK)
        if leftPointer == -1: return 'Not found'
        rightPointer = frontPageText.find(presentation_config.RIGHT_AUTHOR_MARK)
        return frontPageText[leftPointer + len(presentation_config.LEFT_AUTHOR_MARK): rightPointer].strip()
    
    @staticmethod
    def __getGoalAndTasks(slides: Slides) -> str:
        '''Method to find text about goal and tasks. 
        Gets info from slide with title == config.GOAL_AND_TASKS_SLIDE_TITLE.
        Args:
            slides: Slides - presentation slides sequence
        Returns:
            str - goal and tasks text or `Not found`
        '''
        slide = None
        for slide_ in slides:
            if PresentationParserService.__getSlideTitle(slide_).lower() == presentation_config.GOAL_AND_TASKS_SLIDE_TITLE:
                slide = slide_
                break
        if slide is None:
            return 'Not found'
        return PresentationParserService.__getSlideText(slide)

    @staticmethod
    @presentation_utils.deleteSpecialSymbolsFromOutput
    def __getSlideTitle(slide: Slide) -> str:
        '''Method to get title from slide. 
        Args:
            slide: Slide - slide object
        Returns:
            str - title or `Not found`
        '''
        if not slide.shapes or not slide.shapes.title or not slide.shapes.title.text.strip():
            return 'Not found'
        return slide.shapes.title.text
    
    @staticmethod
    @presentation_utils.deleteSpecialSymbolsFromOutput
    def __getSlideText(slide: Slide) -> str:
        '''Method to get full text from slide. 
        Args:
            slide: Slide - slide object
        Returns:
            str - text from slide or empty string
        '''
        if not slide or not slide.shapes:
            return 'Not found'
        text_runs = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        return ' '.join(text_runs)
