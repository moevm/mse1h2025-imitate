from pptx import Presentation
from pptx.slide import Slides, Slide
from typing import IO, Union

from graduate_imitator.apps.graduation.domain.dto.presentation_data import PresentationData
from . import config
from graduate_imitator.apps.graduation.utils import presentation_utils 


class PresentationParser:
    '''Class for parsing PPTX presentations'''

    @staticmethod
    def parsePPTX(file: Union[str, IO[bytes]]) -> PresentationData:
        '''Method to parse all needed info from pptx presentation
        Args:
            file: str or bytes array - path to pptx presentation file or bytes array
        Returns:
            PresentationData object
        Raises:
            pptx.exc.PackageNotFoundError - if file does not exists or file with invalid format
        ''' 
        presentation = Presentation(file)
        topic = PresentationParser.__getTopic(presentation.slides)
        goalAndTasks = PresentationParser.__getGoalAndTasks(presentation.slides)
        author = PresentationParser.__getAuthor(presentation.slides)
        slidesTitles = [PresentationParser.__getSlideTitle(slide) for slide in presentation.slides]
        return PresentationData(
            topic=topic,
            goalAndTasks=goalAndTasks,
            author=author,
            slidesTitles=slidesTitles
        )
    
    @staticmethod
    def __getTopic(slides: Slides) -> str:
        '''Method to find topic of presentation. Just trying to get title at first slide.
        Args:
            slides: Slides - presentation slides sequence
        Returns:
            str - title or `Not found`
        '''
        if not slides:
            return 'Not found'
        return PresentationParser.__getSlideTitle(slides[0])
    
    @staticmethod
    def __getAuthor(slides: Slides) -> str:
        '''Method to find author of presentation. 
        Gets str between two marks in text at first slide (config.LEFT_AUTHOR_MARK, config.RIGHT_AUTHOR_MARK).
        Args:
            slides: Slides - presentation slides sequence
        Returns:
            str - author or `Not found`
        '''
        if not slides:
            return 'Not found'
        frontPageText = PresentationParser.__getSlideText(slides[0]).lower()
        leftPointer = frontPageText.find(config.LEFT_AUTHOR_MARK)
        if leftPointer == -1: return 'Not found'
        rightPointer = frontPageText.find(config.RIGHT_AUTHOR_MARK)
        return frontPageText[leftPointer + len(config.LEFT_AUTHOR_MARK): rightPointer].strip()
    
    @staticmethod
    def __getGoalAndTasks(slides: Slides) -> str:
        '''Method to find text about goal and tasks. 
        Gets info from slide with title == config.GOAL_AND_TASKS_SLIDE_TITLE.
        Args:
            slides: Slides - presentation slides sequence
        Returns:
            str - goal and tasks text or `Not found`
        '''
        slide = None
        for slide_ in slides:
            if PresentationParser.__getSlideTitle(slide_).lower() == config.GOAL_AND_TASKS_SLIDE_TITLE:
                slide = slide_
                break
        if slide is None:
            return 'Not found'
        return PresentationParser.__getSlideText(slide)

    @staticmethod
    @presentation_utils.deleteSpecialSymbolsFromOutput
    def __getSlideTitle(slide: Slide) -> str:
        '''Method to get title from slide. 
        Args:
            slide: Slide - slide object
        Returns:
            str - title or `Not found`
        '''
        if not slide.shapes or not slide.shapes.title or not slide.shapes.title.text.strip():
            return 'Not found'
        return slide.shapes.title.text
    
    @staticmethod
    @presentation_utils.deleteSpecialSymbolsFromOutput
    def __getSlideText(slide: Slide) -> str:
        '''Method to get full text from slide. 
        Args:
            slide: Slide - slide object
        Returns:
            str - text from slide or empty string
        '''
        if not slide or not slide.shapes:
            return 'Not found'
        text_runs = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        return ' '.join(text_runs)