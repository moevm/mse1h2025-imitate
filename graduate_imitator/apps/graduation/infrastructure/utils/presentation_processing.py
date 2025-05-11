from graduate_imitator.apps.graduation.infrastructure.parsers.presentation_parser.presentation_parser import PresentationParser
from graduate_imitator.apps.graduation.utils.presentation_processing_utils import get_10_keywords, clear_text
from pptx import Presentation
from pptx.slide import Slides, Slide
from zipfile import BadZipFile


class PresentationProcessingService:
    @staticmethod
    def get_10_keywords(presentation_file):
        try:
            presentation = Presentation(presentation_file)
        except BadZipFile:
            return ["Bad presentation file (perhaps invalid file format)"]
        slides_text = []

        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(clear_text(run.text))
            slides_text.append(" ".join(slide_text).lower())

        keywords = get_10_keywords(slides_text)
        return keywords
